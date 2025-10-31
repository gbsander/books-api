from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Criar apresentação
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Slide de título"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Título
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 102, 204)
    title_para.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets):
    """Slide com título e bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 102, 204)

    # Conteúdo
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.text = bullet
        p.font.size = Pt(20)
        p.level = 0
        p.space_before = Pt(12)

    return slide

# Slide 1: Título
add_title_slide(prs,
    "Books API",
    "API RESTful com Web Scraping\nTech Challenge - Guilherme Becker Sander")

# Slide 2: Visão Geral
add_content_slide(prs, "Visão Geral do Projeto", [
    "✓ Web Scraping de 1000 livros de books.toscrape.com",
    "✓ API RESTful com FastAPI (5 endpoints)",
    "✓ Documentação completa (Swagger automático)",
    "✓ Pipeline de dados end-to-end",
    "✓ Pronto para deploy em produção"
])

# Slide 3: Arquitetura
add_content_slide(prs, "Arquitetura do Sistema", [
    "📊 Pipeline de Dados:",
    "   1. Ingestão → Web Scraping (BeautifulSoup)",
    "   2. Armazenamento → CSV (1000 livros)",
    "   3. API → FastAPI + Pydantic",
    "   4. Consumo → Cientistas de dados / ML"
])

# Slide 4: Web Scraping
add_content_slide(prs, "Web Scraping", [
    "🔍 Tecnologias:",
    "   • requests - Requisições HTTP",
    "   • BeautifulSoup4 - Parse de HTML",
    "   • pandas - Manipulação de dados",
    "",
    "⚡ Features:",
    "   • Paginação automática (50 páginas)",
    "   • Delays éticos (0.5s entre requests)",
    "   • Extração: título, preço, rating, disponibilidade, imagem"
])

# Slide 5: API RESTful
add_content_slide(prs, "API RESTful - Endpoints", [
    "📡 5 Endpoints Implementados:",
    "",
    "GET /api/v1/health - Status da API",
    "GET /api/v1/books - Lista todos os livros",
    "GET /api/v1/books/{id} - Busca por ID",
    "GET /api/v1/books/search?title=... - Busca por título",
    "GET /api/v1/categories - Lista categorias"
])

# Slide 6: Tecnologias
add_content_slide(prs, "Stack Tecnológica", [
    "🐍 Backend:",
    "   • Python 3.11+",
    "   • FastAPI - Framework web moderno",
    "   • Uvicorn - Servidor ASGI",
    "   • Pydantic - Validação automática",
    "",
    "📦 Deploy:",
    "   • Render / Heroku",
    "   • GitHub para versionamento",
    "   • Documentação Swagger automática"
])

# Slide 7: Escalabilidade
add_content_slide(prs, "Plano de Escalabilidade", [
    "📈 Fase 1 (Atual): MVP",
    "   CSV → API em memória → ~1000 livros",
    "",
    "📈 Fase 2: Banco de Dados",
    "   PostgreSQL → ORM → Filtros complexos",
    "",
    "📈 Fase 3: Cache & Performance",
    "   Redis → 90% redução de latência",
    "",
    "📈 Fase 4: Arquitetura Distribuída",
    "   Load Balancer → Múltiplas instâncias → ~1M livros"
])

# Slide 8: Integração com ML
add_content_slide(prs, "Integração com Machine Learning", [
    "🤖 Casos de Uso:",
    "",
    "Sistema de Recomendação",
    "   • Features: preço, rating, categoria",
    "   • Modelo: K-Nearest Neighbors",
    "",
    "Previsão de Disponibilidade",
    "   • Classificação: In stock vs Out of stock",
    "   • Modelo: Random Forest",
    "",
    "Pipeline MLOps pronto para integração"
])

# Slide 9: Demonstração
add_content_slide(prs, "Demonstração ao Vivo", [
    "💻 Scraping:",
    "   python scripts/scraper.py",
    "",
    "🚀 API Local:",
    "   uvicorn api.main:app --reload",
    "   http://localhost:8000/docs",
    "",
    "🌐 Produção:",
    "   https://books-api.onrender.com",
    "   Swagger interativo disponível"
])

# Slide 10: Documentação
add_content_slide(prs, "Documentação Completa", [
    "📚 Documentos Criados:",
    "",
    "• README.md - Documentação principal",
    "• ARQUITETURA.md - Plano arquitetural detalhado",
    "• DEPLOY.md - Guia de deploy",
    "• ENDPOINTS.md - Referência da API",
    "• RESUMO_APRENDIZADO.md - Conceitos aplicados",
    "",
    "GitHub: github.com/gbsander/books-api"
])

# Slide 11: Resultados
add_content_slide(prs, "Resultados Alcançados", [
    "✅ 1000 livros extraídos e processados",
    "✅ API funcional com 5 endpoints",
    "✅ Documentação Swagger automática",
    "✅ Código versionado no GitHub",
    "✅ Deploy em produção (Render)",
    "✅ Latência < 100ms para consultas",
    "✅ Arquitetura escalável documentada",
    "✅ Pronto para integração com ML"
])

# Slide 12: Próximos Passos
add_content_slide(prs, "Próximos Passos", [
    "🔮 Melhorias Futuras:",
    "",
    "• Migração para PostgreSQL",
    "• Implementação de cache (Redis)",
    "• Autenticação JWT",
    "• Rate limiting",
    "• Paginação nos endpoints",
    "• Testes automatizados (pytest)",
    "• CI/CD com GitHub Actions",
    "• Monitoramento (Prometheus + Grafana)"
])

# Slide 13: Agradecimentos
add_title_slide(prs,
    "Obrigado!",
    "Guilherme Becker Sander\ngithub.com/gbsander/books-api\ngbsander@gmail.com")

# Salvar apresentação
prs.save('TechChallenge_BooksAPI.pptx')
print("✅ Apresentação criada: TechChallenge_BooksAPI.pptx")
print("📊 Total de slides: 13")
